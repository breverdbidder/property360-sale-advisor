"""
Generate 07_offering_memorandum.pptx for Palm Bay Palms Apartments case study.
10-slide institutional-quality Offering Memorandum in 16:9 widescreen.
"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from data import *

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CLR_NAVY = RGBColor(0x1E, 0x3A, 0x5F)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BLACK = RGBColor(0x00, 0x00, 0x00)
CLR_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
CLR_MED_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
CLR_GOLD = RGBColor(0xC5, 0x9A, 0x2C)
CLR_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

FONT_NAME = "Calibri"

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=CLR_BLACK, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, bold=False,
                  color=CLR_DARK_TEXT, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(4), space_after=Pt(2), bullet=False):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    if bullet:
        p.level = 0
    return p


def add_slide_number(slide, slide_num, total=10):
    """Add slide number in bottom-right footer area."""
    add_textbox(
        slide,
        left=SLIDE_WIDTH - Inches(1.5),
        top=SLIDE_HEIGHT - Inches(0.4),
        width=Inches(1.3),
        height=Inches(0.3),
        text=f"{slide_num} / {total}",
        font_size=9,
        color=CLR_MED_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


def add_title_bar(slide, title_text, slide_num):
    """Add a navy title bar across the top of a content slide."""
    # Navy rectangle at top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_NAVY
    shape.line.fill.background()

    # Title text
    add_textbox(
        slide,
        left=Inches(0.6), top=Inches(0.15),
        width=Inches(11), height=Inches(0.8),
        text=title_text,
        font_size=28, bold=True, color=CLR_WHITE,
    )

    # Slide number
    add_slide_number(slide, slide_num)


def add_table(slide, rows, cols, left, top, width, height):
    """Add a table shape and return the table object."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_table_header(table, col_count, font_size=11):
    """Style the first row of a table as navy header."""
    for col_idx in range(col_count):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = CLR_WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_table_body(table, row_count, col_count, font_size=10):
    """Style body rows with alternating shading."""
    for row_idx in range(1, row_count):
        for col_idx in range(col_count):
            cell = table.cell(row_idx, col_idx)
            # Alternating row fill
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = CLR_DARK_TEXT
                paragraph.font.name = FONT_NAME
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_total_row(table, row_idx, col_count, font_size=10):
    """Style a totals row with bold text and navy-ish background."""
    for col_idx in range(col_count):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_MED_GRAY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = CLR_NAVY
            paragraph.font.name = FONT_NAME
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_text(table, row, col, text, alignment=PP_ALIGN.CENTER):
    """Set text for a table cell."""
    cell = table.cell(row, col)
    cell.text_frame.clear()
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.alignment = alignment
    return cell


def set_col_widths(table, widths_inches):
    """Set column widths from a list of inch values."""
    for i, w in enumerate(widths_inches):
        table.columns[i].width = Inches(w)


# ===========================================================================
# Slide builders
# ===========================================================================

def build_slide_01_cover(prs):
    """Slide 1 -- Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, CLR_NAVY)

    # Gold accent line at top
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(1.2),
        Inches(9.333), Inches(0.04),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = CLR_GOLD
    line_shape.line.fill.background()

    # Property name
    add_textbox(
        slide,
        left=Inches(2), top=Inches(1.5),
        width=Inches(9.333), height=Inches(1.0),
        text=PROPERTY["name"],
        font_size=40, bold=True, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    add_textbox(
        slide,
        left=Inches(2), top=Inches(2.5),
        width=Inches(9.333), height=Inches(0.6),
        text="18-Unit Value-Add Multifamily Investment Opportunity",
        font_size=20, bold=False, color=CLR_GOLD,
        alignment=PP_ALIGN.CENTER,
    )

    # Price
    add_textbox(
        slide,
        left=Inches(2), top=Inches(3.4),
        width=Inches(9.333), height=Inches(0.8),
        text=f"${ASKING_PRICE:,.0f}",
        font_size=36, bold=True, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Gold accent line below price
    line_shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4.3),
        Inches(5.333), Inches(0.04),
    )
    line_shape2.fill.solid()
    line_shape2.fill.fore_color.rgb = CLR_GOLD
    line_shape2.line.fill.background()

    # Location
    add_textbox(
        slide,
        left=Inches(2), top=Inches(4.6),
        width=Inches(9.333), height=Inches(0.5),
        text="Palm Bay, Florida",
        font_size=18, bold=False, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Branding
    add_textbox(
        slide,
        left=Inches(2), top=Inches(5.5),
        width=Inches(9.333), height=Inches(0.5),
        text="Presented by Property360 Real Estate",
        font_size=16, bold=False, color=CLR_MED_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Slide number
    add_textbox(
        slide,
        left=SLIDE_WIDTH - Inches(1.5),
        top=SLIDE_HEIGHT - Inches(0.4),
        width=Inches(1.3), height=Inches(0.3),
        text="1 / 10",
        font_size=9, color=CLR_MED_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


def build_slide_02_highlights(prs):
    """Slide 2 -- Investment Highlights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Investment Highlights", 2)

    highlights = [
        f"Below-market rents with $3,800/month ($45,600/yr) upside through lease renewals",
        f"Manageable CapEx profile \u2014 ${TOTAL_CAPEX:,.0f} total, with only $16,000 immediate",
        "Strong Space Coast demographics \u2014 population growth, employment diversification",
        "1031 Exchange eligible \u2014 clean title, single LLC ownership",
    ]

    y_start = Inches(1.6)
    for i, text in enumerate(highlights):
        # Bullet number box (navy circle-ish)
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y_start + Inches(i * 1.2),
            Inches(0.5), Inches(0.5),
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = CLR_NAVY
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CLR_WHITE
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Highlight text
        add_textbox(
            slide,
            left=Inches(1.6), top=y_start + Inches(i * 1.2) + Inches(0.05),
            width=Inches(10.5), height=Inches(0.5),
            text=text,
            font_size=18, bold=False, color=CLR_DARK_TEXT,
        )


def build_slide_03_overview(prs):
    """Slide 3 -- Property Overview (two-column table layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Property Overview", 3)

    # Left column data
    left_data = [
        ("Address", PROPERTY["address"]),
        ("Property Type", PROPERTY["property_type"]),
        ("Year Built", str(PROPERTY["year_built"])),
        ("Total SF", f"{PROPERTY['total_sqft']:,}"),
    ]

    # Right column data
    right_data = [
        ("Lot Size", f"{PROPERTY['lot_acres']} acres"),
        ("Zoning", PROPERTY["zoning"].split(" (")[0]),  # "RM-13"
        ("Flood Zone", PROPERTY["flood_zone"]),
        ("Parking", "Surface lot, 36 spaces"),
    ]

    # Left table
    tbl_left = add_table(slide, len(left_data), 2,
                         Inches(0.6), Inches(1.5),
                         Inches(5.8), Inches(3.5))
    set_col_widths(tbl_left, [2.0, 3.8])

    for r, (label, value) in enumerate(left_data):
        set_cell_text(tbl_left, r, 0, label, PP_ALIGN.LEFT)
        set_cell_text(tbl_left, r, 1, value, PP_ALIGN.LEFT)
        # Style label column
        for p in tbl_left.cell(r, 0).text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = CLR_NAVY
            p.font.name = FONT_NAME
        for p in tbl_left.cell(r, 1).text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = CLR_DARK_TEXT
            p.font.name = FONT_NAME
        # Alternating fill
        fill_color = CLR_LIGHT_GRAY if r % 2 == 0 else CLR_WHITE
        tbl_left.cell(r, 0).fill.solid()
        tbl_left.cell(r, 0).fill.fore_color.rgb = fill_color
        tbl_left.cell(r, 1).fill.solid()
        tbl_left.cell(r, 1).fill.fore_color.rgb = fill_color

    # Right table
    tbl_right = add_table(slide, len(right_data), 2,
                          Inches(6.9), Inches(1.5),
                          Inches(5.8), Inches(3.5))
    set_col_widths(tbl_right, [2.0, 3.8])

    for r, (label, value) in enumerate(right_data):
        set_cell_text(tbl_right, r, 0, label, PP_ALIGN.LEFT)
        set_cell_text(tbl_right, r, 1, value, PP_ALIGN.LEFT)
        for p in tbl_right.cell(r, 0).text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = CLR_NAVY
            p.font.name = FONT_NAME
        for p in tbl_right.cell(r, 1).text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = CLR_DARK_TEXT
            p.font.name = FONT_NAME
        fill_color = CLR_LIGHT_GRAY if r % 2 == 0 else CLR_WHITE
        tbl_right.cell(r, 0).fill.solid()
        tbl_right.cell(r, 0).fill.fore_color.rgb = fill_color
        tbl_right.cell(r, 1).fill.solid()
        tbl_right.cell(r, 1).fill.fore_color.rgb = fill_color

    # Owner info note
    add_textbox(
        slide,
        left=Inches(0.6), top=Inches(5.5),
        width=Inches(12), height=Inches(0.5),
        text=f"Owner: {PROPERTY['owner_entity']}  |  Acquired: {PROPERTY['purchase_date']}  |  Purchase Price: ${PROPERTY['purchase_price']:,.0f}",
        font_size=12, color=CLR_MED_GRAY, alignment=PP_ALIGN.LEFT,
    )


def build_slide_04_unit_mix(prs):
    """Slide 4 -- Unit Mix & Rent Schedule."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Unit Mix & Rent Schedule", 4)

    # Compute unit mix stats from data
    unit_types = {}
    for u in UNITS:
        utype = u[U_TYPE]
        if utype not in unit_types:
            unit_types[utype] = {"count": 0, "sf": u[U_SF], "rents": [],
                                 "market": u[U_MARKET]}
        unit_types[utype]["count"] += 1
        if u[U_STATUS] == "Occupied":
            unit_types[utype]["rents"].append(u[U_RENT])

    headers = ["Type", "Count", "Avg SF", "Current Avg Rent", "Market Rent",
               "Monthly Upside"]
    rows_data = []
    total_upside = 0

    for utype in ["1BR/1BA", "2BR/1BA", "3BR/2BA"]:
        info = unit_types[utype]
        avg_rent = (sum(info["rents"]) / len(info["rents"])
                    if info["rents"] else 0)
        monthly_upside = info["market"] * info["count"] - sum(info["rents"])
        # Account for vacant units: their upside is included as market rent
        vacant_count = info["count"] - len(info["rents"])
        total_upside += monthly_upside

        rent_display = f"${avg_rent:,.0f}*" if vacant_count > 0 else f"${avg_rent:,.0f}"
        rows_data.append([
            utype,
            str(info["count"]),
            f"{info['sf']:,}",
            rent_display,
            f"${info['market']:,.0f}",
            f"${monthly_upside:,.0f}",
        ])

    # Total row
    rows_data.append([
        "Total", str(TOTAL_UNITS), "", "", "",
        f"${total_upside:,.0f}",
    ])

    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)

    tbl = add_table(slide, num_rows, num_cols,
                    Inches(0.8), Inches(1.5),
                    Inches(11.7), Inches(3.5))

    col_widths = [1.6, 1.0, 1.2, 2.4, 2.0, 2.2]
    set_col_widths(tbl, col_widths)

    # Header
    for c, h in enumerate(headers):
        set_cell_text(tbl, 0, c, h, PP_ALIGN.CENTER)
    style_table_header(tbl, num_cols, font_size=12)

    # Body rows
    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            set_cell_text(tbl, r + 1, c, val, PP_ALIGN.CENTER)

    style_table_body(tbl, num_rows, num_cols, font_size=12)

    # Style total row
    style_total_row(tbl, num_rows - 1, num_cols, font_size=12)

    # Footnote
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(5.3),
        width=Inches(10), height=Inches(0.4),
        text="*Average of occupied units only",
        font_size=10, color=CLR_MED_GRAY, alignment=PP_ALIGN.LEFT,
    )


def build_slide_05_financial(prs):
    """Slide 5 -- Financial Performance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Financial Performance", 5)

    # Metrics in 2x3 card layout
    metrics = [
        ("Trailing 12 NOI", f"${NOI_ACTUAL:,.0f}",
         f"Cap Rate: {CAP_RATE_ACTUAL:.2%}"),
        ("Pro Forma NOI", f"${NOI_PROFORMA:,.0f}",
         f"Cap Rate: {CAP_RATE_PROFORMA:.2%}"),
        ("Total Revenue (T12)", f"${TOTAL_REVENUE_ACTUAL:,.0f}", ""),
        ("Expense Ratio",
         f"{TOTAL_EXPENSES_ACTUAL / TOTAL_REVENUE_ACTUAL:.1%} (T12)",
         f"\u2192 {TOTAL_EXPENSES_PROFORMA / TOTAL_REVENUE_PROFORMA:.1%} (stabilized)"),
        ("Cash Flow After DS", f"${CASH_FLOW_AFTER_DS:,.0f}/yr", ""),
        ("Annual Debt Service", f"${ANNUAL_DEBT_SERVICE:,.0f}",
         f"Mortgage Rate: {PROPERTY['mortgage_rate']:.2%}"),
    ]

    positions = [
        (Inches(0.6), Inches(1.5)),   (Inches(4.7), Inches(1.5)),
        (Inches(8.8), Inches(1.5)),
        (Inches(0.6), Inches(4.0)),   (Inches(4.7), Inches(4.0)),
        (Inches(8.8), Inches(4.0)),
    ]

    card_w = Inches(3.8)
    card_h = Inches(2.0)

    for i, (label, value, subtitle) in enumerate(metrics):
        left, top = positions[i]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_w, card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CLR_LIGHT_GRAY
        card.line.color.rgb = CLR_MED_GRAY
        card.line.width = Pt(0.5)

        # Label
        add_textbox(
            slide,
            left=left + Inches(0.2), top=top + Inches(0.15),
            width=card_w - Inches(0.4), height=Inches(0.4),
            text=label,
            font_size=12, bold=True, color=CLR_NAVY,
        )

        # Value
        add_textbox(
            slide,
            left=left + Inches(0.2), top=top + Inches(0.6),
            width=card_w - Inches(0.4), height=Inches(0.6),
            text=value,
            font_size=24, bold=True, color=CLR_DARK_TEXT,
        )

        # Subtitle
        if subtitle:
            add_textbox(
                slide,
                left=left + Inches(0.2), top=top + Inches(1.3),
                width=card_w - Inches(0.4), height=Inches(0.4),
                text=subtitle,
                font_size=11, bold=False, color=CLR_MED_GRAY,
            )


def build_slide_06_rent_comps(prs):
    """Slide 6 -- Market Rent Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Market Rent Analysis", 6)

    headers = ["Unit Type", "Current Avg Rent", "Market Rent", "Upside/Unit",
               "Upside/Month (All)"]

    # Recompute from data
    unit_types_info = {}
    for u in UNITS:
        utype = u[U_TYPE]
        if utype not in unit_types_info:
            unit_types_info[utype] = {"count": 0, "occ_rents": [],
                                      "market": u[U_MARKET]}
        unit_types_info[utype]["count"] += 1
        if u[U_STATUS] == "Occupied":
            unit_types_info[utype]["occ_rents"].append(u[U_RENT])

    rows_data = []
    for utype in ["1BR/1BA", "2BR/1BA", "3BR/2BA"]:
        info = unit_types_info[utype]
        avg_rent = (sum(info["occ_rents"]) / len(info["occ_rents"])
                    if info["occ_rents"] else 0)
        upside_per_unit = info["market"] - avg_rent
        total_upside = upside_per_unit * len(info["occ_rents"])
        rows_data.append([
            utype,
            f"${avg_rent:,.0f}",
            f"${info['market']:,.0f}",
            f"${upside_per_unit:,.0f}",
            f"${total_upside:,.0f}",
        ])

    num_rows = len(rows_data) + 1
    num_cols = len(headers)
    tbl = add_table(slide, num_rows, num_cols,
                    Inches(0.8), Inches(1.5),
                    Inches(11.7), Inches(2.5))
    set_col_widths(tbl, [2.0, 2.5, 2.2, 2.2, 2.8])

    for c, h in enumerate(headers):
        set_cell_text(tbl, 0, c, h, PP_ALIGN.CENTER)
    style_table_header(tbl, num_cols, font_size=12)

    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            set_cell_text(tbl, r + 1, c, val, PP_ALIGN.CENTER)
    style_table_body(tbl, num_rows, num_cols, font_size=12)

    # Market context note
    txBox = add_textbox(
        slide,
        left=Inches(0.8), top=Inches(4.5),
        width=Inches(11), height=Inches(1.5),
        text="Market Context:",
        font_size=14, bold=True, color=CLR_NAVY,
    )
    tf = txBox.text_frame
    add_paragraph(
        tf,
        "Palm Bay 2BR market average: $1,400 - $1,500/month (per CoStar, Zillow, and Apartments.com data)",
        font_size=12, color=CLR_DARK_TEXT,
    )
    add_paragraph(
        tf,
        "Subject property's below-market units present significant rent upside upon lease expiration or renewal",
        font_size=12, color=CLR_DARK_TEXT,
    )
    add_paragraph(
        tf,
        "Recent renovated comps in the area achieving $1,450 - $1,550 for 2BR units",
        font_size=12, color=CLR_DARK_TEXT,
    )


def build_slide_07_capex(prs):
    """Slide 7 -- Capital Expenditure Plan."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Capital Expenditure Plan", 7)

    headers = ["Item", "Priority", "Cost", "Timeline"]

    # Combine unit turns into a single line per spec
    capex_display = [
        ("Roof replacement", "Medium", "$50,000", "3-5 years"),
        ("HVAC (3 units)", "High", "$13,500", "Immediate"),
        ("Water heaters (6)", "Medium", "$4,800", "1-2 years"),
        ("Stucco repair", "Low", "$8,000", "1 year"),
        ("Unit turns (2)", "High", "$5,000", "Before listing"),
        ("Parking lot reseal", "Low", "$6,000", "1 year"),
    ]

    num_rows = len(capex_display) + 2  # header + body + total
    num_cols = len(headers)
    tbl = add_table(slide, num_rows, num_cols,
                    Inches(1.5), Inches(1.5),
                    Inches(10.3), Inches(4.5))
    set_col_widths(tbl, [3.5, 1.8, 2.0, 3.0])

    # Header
    for c, h in enumerate(headers):
        set_cell_text(tbl, 0, c, h, PP_ALIGN.CENTER)
    style_table_header(tbl, num_cols, font_size=12)

    # Body
    for r, (item, priority, cost, timeline) in enumerate(capex_display):
        set_cell_text(tbl, r + 1, 0, item, PP_ALIGN.LEFT)
        set_cell_text(tbl, r + 1, 1, priority, PP_ALIGN.CENTER)
        set_cell_text(tbl, r + 1, 2, cost, PP_ALIGN.CENTER)
        set_cell_text(tbl, r + 1, 3, timeline, PP_ALIGN.CENTER)

    # Total row
    total_row_idx = num_rows - 1
    set_cell_text(tbl, total_row_idx, 0, "Total", PP_ALIGN.LEFT)
    set_cell_text(tbl, total_row_idx, 1, "", PP_ALIGN.CENTER)
    set_cell_text(tbl, total_row_idx, 2, f"${TOTAL_CAPEX:,.0f}", PP_ALIGN.CENTER)
    set_cell_text(tbl, total_row_idx, 3, "", PP_ALIGN.CENTER)

    style_table_body(tbl, num_rows, num_cols, font_size=12)
    style_total_row(tbl, total_row_idx, num_cols, font_size=12)


def build_slide_08_demographics(prs):
    """Slide 8 -- Palm Bay, FL Market Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Palm Bay, FL \u2014 Market Overview", 8)

    stats = [
        ("Population", "~125,000 (10% growth 2020-2025)"),
        ("Median Household Income", "~$55,000"),
        ("Major Employers",
         "L3Harris Technologies, Northrop Grumman, Health First"),
        ("Location",
         "Florida's Space Coast, 70 miles east of Orlando"),
        ("Brevard County Vacancy Rate", "~5%"),
    ]

    y_start = Inches(1.6)
    for i, (label, value) in enumerate(stats):
        row_top = y_start + Inches(i * 1.05)

        # Background stripe
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), row_top,
            Inches(12.1), Inches(0.85),
        )
        stripe_color = CLR_LIGHT_GRAY if i % 2 == 0 else CLR_WHITE
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = stripe_color
        stripe.line.fill.background()

        # Label
        add_textbox(
            slide,
            left=Inches(0.9), top=row_top + Inches(0.15),
            width=Inches(3.5), height=Inches(0.55),
            text=label,
            font_size=16, bold=True, color=CLR_NAVY,
        )

        # Value
        add_textbox(
            slide,
            left=Inches(4.8), top=row_top + Inches(0.15),
            width=Inches(7.5), height=Inches(0.55),
            text=value,
            font_size=16, bold=False, color=CLR_DARK_TEXT,
        )


def build_slide_09_proforma(prs):
    """Slide 9 -- 3-Year Pro Forma Projections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "3-Year Pro Forma", 9)

    # Assumptions
    add_textbox(
        slide,
        left=Inches(0.6), top=Inches(1.2),
        width=Inches(12), height=Inches(0.4),
        text="Assumptions: 3% annual rent growth  |  2% annual expense growth  |  Vacancy: 11.1% (Y1) \u2192 7% (Y2) \u2192 5% (Y3)",
        font_size=11, bold=False, color=CLR_MED_GRAY,
    )

    # Year 1: current performance (T12 actuals)
    # Year 2: partial stabilization (rent growth + lower vacancy)
    # Year 3: stabilized (market rents + 5% vacancy)

    # Year 1 numbers (T12 actuals)
    y1_revenue = TOTAL_REVENUE_ACTUAL   # $258,000
    y1_expenses = TOTAL_EXPENSES_ACTUAL  # $123,750
    y1_noi = NOI_ACTUAL                  # $134,250

    # Year 2: GPR grows 3%, vacancy drops to 7%, expenses grow 2%
    y2_gpr = GPR_ACTUAL * 1.03           # $290,460
    y2_vacancy = y2_gpr * 0.07           # ~$20,332
    y2_egi = y2_gpr - y2_vacancy
    y2_other = (LAUNDRY_INCOME_ACTUAL + LATE_FEES_ACTUAL) * 1.03
    y2_revenue = y2_egi + y2_other
    y2_expenses_base = sum(
        e[1] for e in EXPENSES_ACTUAL if e[1] is not None
    ) * 1.02
    y2_mgmt = y2_revenue * MGMT_FEE_PCT
    y2_expenses = y2_expenses_base + y2_mgmt
    y2_noi = y2_revenue - y2_expenses

    # Year 3: GPR grows another 3% from Y2, vacancy 5%, expenses grow 2%
    y3_gpr = y2_gpr * 1.03
    y3_vacancy = y3_gpr * 0.05
    y3_egi = y3_gpr - y3_vacancy
    y3_other = y2_other * 1.03
    y3_revenue = y3_egi + y3_other
    y3_expenses_base = y2_expenses_base * 1.02
    y3_mgmt = y3_revenue * MGMT_FEE_PCT
    y3_expenses = y3_expenses_base + y3_mgmt
    y3_noi = y3_revenue - y3_expenses

    headers = ["", "Year 1 (T12)", "Year 2", "Year 3"]
    rows_data = [
        ("Gross Potential Rent", f"${GPR_ACTUAL:,.0f}",
         f"${y2_gpr:,.0f}", f"${y3_gpr:,.0f}"),
        ("Less: Vacancy", f"$({VACANCY_LOSS_ACTUAL_PRD:,.0f})",
         f"$({y2_vacancy:,.0f})", f"$({y3_vacancy:,.0f})"),
        ("Total Revenue", f"${y1_revenue:,.0f}",
         f"${y2_revenue:,.0f}", f"${y3_revenue:,.0f}"),
        ("Total Expenses", f"$({y1_expenses:,.0f})",
         f"$({y2_expenses:,.0f})", f"$({y3_expenses:,.0f})"),
        ("Net Operating Income", f"${y1_noi:,.0f}",
         f"${y2_noi:,.0f}", f"${y3_noi:,.0f}"),
        ("Cash Flow After DS", f"${y1_noi - ANNUAL_DEBT_SERVICE:,.0f}",
         f"${y2_noi - ANNUAL_DEBT_SERVICE:,.0f}",
         f"${y3_noi - ANNUAL_DEBT_SERVICE:,.0f}"),
    ]

    num_rows = len(rows_data) + 1
    num_cols = len(headers)
    tbl = add_table(slide, num_rows, num_cols,
                    Inches(1.0), Inches(1.8),
                    Inches(11.3), Inches(4.5))
    set_col_widths(tbl, [3.5, 2.6, 2.6, 2.6])

    for c, h in enumerate(headers):
        set_cell_text(tbl, 0, c, h, PP_ALIGN.CENTER)
    style_table_header(tbl, num_cols, font_size=12)

    for r, row_data in enumerate(rows_data):
        set_cell_text(tbl, r + 1, 0, row_data[0], PP_ALIGN.LEFT)
        for c in range(1, num_cols):
            set_cell_text(tbl, r + 1, c, row_data[c], PP_ALIGN.CENTER)

    style_table_body(tbl, num_rows, num_cols, font_size=12)

    # Bold the NOI and Cash Flow rows
    for highlight_row in [num_rows - 2, num_rows - 1]:  # NOI and CF rows
        for c in range(num_cols):
            cell = tbl.cell(highlight_row, c)
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = CLR_NAVY


def build_slide_10_terms(prs):
    """Slide 10 -- Transaction Summary / Offer Terms."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CLR_WHITE)
    add_title_bar(slide, "Transaction Summary", 10)

    terms = [
        ("Asking Price", f"${ASKING_PRICE:,.0f}"),
        ("Price Per Unit", f"${ASKING_PRICE // TOTAL_UNITS:,.0f}"),
        ("Earnest Money", "$50,000 (hard after 15-day inspection)"),
        ("Inspection Period", "15 business days"),
        ("Financing Contingency", "30 days"),
        ("Target Close", "45 days from execution"),
    ]

    y_start = Inches(1.5)
    for i, (label, value) in enumerate(terms):
        row_top = y_start + Inches(i * 0.85)

        # Stripe
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), row_top,
            Inches(10.3), Inches(0.7),
        )
        stripe_color = CLR_LIGHT_GRAY if i % 2 == 0 else CLR_WHITE
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = stripe_color
        stripe.line.fill.background()

        # Label
        add_textbox(
            slide,
            left=Inches(1.8), top=row_top + Inches(0.1),
            width=Inches(3.5), height=Inches(0.5),
            text=label,
            font_size=16, bold=True, color=CLR_NAVY,
        )

        # Value
        add_textbox(
            slide,
            left=Inches(5.5), top=row_top + Inches(0.1),
            width=Inches(6), height=Inches(0.5),
            text=value,
            font_size=16, bold=False, color=CLR_DARK_TEXT,
        )

    # Contact info at bottom
    contact_y = y_start + Inches(len(terms) * 0.85) + Inches(0.5)

    # Navy box for contact
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), contact_y,
        Inches(8.3), Inches(1.0),
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = CLR_NAVY
    contact_box.line.fill.background()

    add_textbox(
        slide,
        left=Inches(2.7), top=contact_y + Inches(0.1),
        width=Inches(7.9), height=Inches(0.4),
        text="Contact",
        font_size=14, bold=True, color=CLR_GOLD,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        left=Inches(2.7), top=contact_y + Inches(0.45),
        width=Inches(7.9), height=Inches(0.4),
        text="Mariam Shapira  |  Property360 Real Estate",
        font_size=16, bold=True, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )


# ===========================================================================
# Main
# ===========================================================================
def main():
    prs = Presentation()

    # Set 16:9 widescreen dimensions
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_01_cover(prs)
    build_slide_02_highlights(prs)
    build_slide_03_overview(prs)
    build_slide_04_unit_mix(prs)
    build_slide_05_financial(prs)
    build_slide_06_rent_comps(prs)
    build_slide_07_capex(prs)
    build_slide_08_demographics(prs)
    build_slide_09_proforma(prs)
    build_slide_10_terms(prs)

    filepath = output_path("07_offering_memorandum.pptx")
    prs.save(filepath)

    slide_count = len(prs.slides)
    print(f"Created 07_offering_memorandum.pptx at {filepath}")
    print(f"Total slides: {slide_count}")

    if slide_count != 10:
        print(f"WARNING: Expected 10 slides, got {slide_count}")
        sys.exit(1)


if __name__ == "__main__":
    main()
